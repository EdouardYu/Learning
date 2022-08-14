from hello.pptx_function import *
from numpy import array
from pptx import Presentation

def hello():
    prs = Presentation()

    title_slide(prs, "Hello World!", "Edouard Yu")

    idées = ["Société De Livraison Des Ouvrages Olympiques",
            "18 rue de Londres",
            "75 009 Paris",
            "Gare Saint-Lazare",
            "Métro 14"]
    niveaux = array([0, 1, 1, 0, 1])
    summary_slide(prs, "SOLIDEO", idées, niveaux)

    picture_slide(prs, "Python")
    
    graph  = graph_slide(prs,"Graphique", "Echantillon", [1, "A", 3,], array([5, 2, 3]))
    category_axis = graph.chart.category_axis
    category_axis.has_major_gridlines = True

    property = ['Prénom', 'Nom', 'Age', 'Passion']
    person1 = ['Edouard', 'Yu', 21, 'Programmation']
    person2 = ['Xinyi', 'Xian', 22, 'Finance']
    data = [property, person1, person2]
    table_slide(prs, "Tableau", data)

    link_slide(prs, "Lien SOLIDEO", "SOLIDEO", "https://www.ouvrages-olympiques.fr/fr")

    prs.save("hello.pptx")

    print("PowerPoint créé/modifié !\nVérifiez qu'il n'y est pas d'erreur")

if __name__ == "__main__":
    hello()


"""
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
# Pour mettre des formes
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
# Pour les graphes
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_TICK_MARK


def hello():
    prs = Presentation()
    # slide_layouts permet de définir la référence de chaque future diapositive
    # 0 =  Title Slide, 1 = Title and Content, 3 = Section Header, etc.
    # (cf PowerPoint -> Nouvelle diapositive -> Office Theme)
    slide_layout1 = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(slide_layout1)
    # Par défaut, shapes.title est le main top placeholder (cadre)
    title1 = slide1.shapes.title
    # Le deuxième placeholder -> 1
    subtitle = slide1.placeholders[1]
    # -------------------------------
    slide2_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(slide2_layout)
    title2 = slide2.shapes.title
    bullet_point_box = slide2.shapes
    bullet_point_boxlvl1 = bullet_point_box.placeholders[1]
    # Pour faire des points de sommaires
    bullet_point_boxlvl1.text = 'Société De Livraison Des Ouvrages Olympiques'
    # Pour faire des sous-points du point de 1 (sous-paragraphes)
    # Remarque : on ne peut pas mettre le texte autre part
    bullet_point_boxlvl2 = bullet_point_boxlvl1.text_frame.add_paragraph()
    bullet_point_boxlvl2.text = '18 rue de Londres'
    # Pour donner le sous-point du point 1:
    # level 1 est le sous-point, le level 2 est le sous-sous-point
    # équivalent du sous-point du sous-point du point 1
    bullet_point_boxlvl2.level = 1
    bullet_point_boxlvl3 = bullet_point_boxlvl1.text_frame.add_paragraph()
    bullet_point_boxlvl3.text = '75 009 Paris'
    bullet_point_boxlvl3.level = 2
    bullet_point_boxlvl1_bis = bullet_point_boxlvl1.text_frame.add_paragraph()
    bullet_point_boxlvl1_bis.text = 'Gare Saint-Lazare'
    bullet_point_boxlvl2_bis = bullet_point_boxlvl1.text_frame.add_paragraph()
    bullet_point_boxlvl2_bis.text = 'Métro 14'
    bullet_point_boxlvl2_bis.level = 1
    # ----------------------------------
    slide3_layout = prs.slide_layouts[5]
    slide3 = prs.slides.add_slide(slide3_layout)
    title3 = slide3.shapes.title
    img = './images/python.png'
    from_left = Inches(3)
    from_top = Inches(2)
    slide3.shapes.add_picture(img, from_left, from_top)
    # Il faut cependant s'assurer de la taille de l'image en amont
    # ----------------------------------
    slide4_layout = prs.slide_layouts[5]
    slide4 = prs.slides.add_slide(slide4_layout)
    title4 = slide4.shapes.title
    left1 = top1 = width1 = height1 = Inches(2)
    # On fait un rectangle aux bords arrondis
    slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            left1, top1, width1, height1)
    left2 = Inches(6)
    top2 = width2 = height2 = Inches(2)
    arrow1 = slide4.shapes.add_shape(MSO_SHAPE.UP_ARROW,
                                     left2, top2, width2, height2)
    # Pour donner une autre couleur
    fill_arrow1 = arrow1.fill
    # Permet de rendre la forme en solide
    fill_arrow1.solid()
    # ACCENT_1 = bleu 2 = rouge 3 = vert etc.
    fill_arrow1.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_3

    arrow1.rotation = -45  # degré sens horaire
    # --------------------------------
    slide5_layout = prs.slide_layouts[5]
    slide5 = prs.slides.add_slide(slide5_layout)
    title5 = slide5.shapes.title
    # Construction d'un graphe
    graph_info = CategoryChartData()
    # Construction des éléments en abscisse
    graph_info.categories = ['A', 'B', 'C']
    # Construction des valeurs en ordonnée
    graph_info.add_series('Echantillon', [15, 11, 8])
    left_graph = top_graph = Inches(2)
    width_graph = Inches(6)
    height_graph = Inches(4)
    # Affichage du graphique
    graph_frame = slide5.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left_graph,
                                          top_graph, width_graph, height_graph, graph_info)
    # Pour manipuler la grille du graphique
    graph = graph_frame.chart
    category_axis = graph.category_axis
    # Pour faire des grilles verticales
    category_axis.has_major_gridlines = True
    # Pour avoir une marque sur les valeurs exactes en abscisse
    # (OUTSIDE met la marque en dehord du cadre et INSIDE à l'intérieur)
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    # Les labels de l'axe des abscisses sont en italique et de taille 24px
    category_axis.tick_labels.font.italic = True
    category_axis.tick_labels.font.size = Pt(24)
    # --------------------------------
    slide6_layout = prs.slide_layouts[5]
    slide6 = prs.slides.add_slide(slide6_layout)
    title6 = slide6.shapes.title
    left_table = Inches(1)
    top_table = Inches(2)
    width_table = Inches(8)
    height_table = Inches(3)
    rows = -3
    columns = 4
    # Création d'un tableau à 3 lignes et 4 colonnes
    table_frame = slide6.shapes.add_table(rows, columns, left_table, top_table,
                                          width_table, height_table)
    table = table_frame.table
    for i in range(columns):
        for j in range(rows):
            title = ['Prénom', 'Nom', 'Age', 'Passion']
            person1 = ['Edouard', 'Yu', '21', 'Programmation']
            person2 = ['Xinyi', 'Xian', '22', 'Finance']
            data = [title, person1, person2]
            cell = table.cell(j, i)
            cell.text = data[j][i]
    # --------------------------------
    slide7_layout = prs.slide_layouts[0]
    slide7 = prs.slides.add_slide(slide7_layout)
    title7 = slide7.shapes.title
    para = slide7.placeholders[1].text_frame.paragraphs[0]
    # Pour créer un lien URL
    addrun = para.add_run()
    addrun.text = 'SOLIDEO'
    hlink = addrun.hyperlink
    hlink.address = 'https://www.ouvrages-olympiques.fr/fr'

    title1.text = 'Hello World!'
    subtitle.text = 'Edouard Y'

    title2.text = 'SOLIDEO'

    title3.text = 'Python'

    title4.text = 'Formes'

    title5.text = 'Graphique'

    title6.text = 'Tableau'

    title7.text = 'Lien URL'

    prs.save('hello.pptx')

    print('PowerPoint créé/modifié !')


if __name__ == "__main__":
    hello()
"""