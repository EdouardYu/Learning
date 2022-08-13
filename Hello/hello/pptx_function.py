from numpy import zeros, ones, array
import collections.abc
from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


def title_slide(presentation, title='', subtitle=''):
    try:
        if not type(presentation) is type(Presentation()):
            raise TypeError
    except TypeError:
        print("TypeError: {} doit être un objet Presentation()\n".format(presentation))
    else:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = str(title)
        slide.shapes.placeholders[1].text =str(subtitle)


def summary_slide(presentation, title = '', bullet_points = [], levels = None):
    if(not type(levels) is type(array([]))):
        levels = zeros(len(bullet_points), int)
        print("Attention, vous n'avez pas renseigner les niveaux sous forme d'un un numpy.array. Vérifiez bien\n")
    try:
        error = None
        if not type(presentation) is type(Presentation()):
            error = 1
            raise TypeError
        if not type(bullet_points) is list:
            error = 2
            raise TypeError
        if len(bullet_points) != len(levels):
            raise ValueError
        if(len(bullet_points) > 8):
            error = 1
            raise AssertionError
        for bullet_point in bullet_points:
            if len(bullet_point) > 44:
                error = 2
                raise AssertionError
    except TypeError:
        if error == 1:
            print("TypeError : {} doit être un objet Presentation()\n".format(presentation))
        elif error == 2:
            print("TypeError : Veuillez renseigner une liste contenant chaques idées du sommaire")
            print("['première_idée', 'autre_idée']\n")
    except ValueError:
        print("ValueError : Les listes des idées et des niveaux doivent être de même taille\n")
    except AssertionError:
        if error == 1:
            print("AssertionError : Vous ne pouvez que rentrer que 8 idées au maximum\n")
        elif error == 2:
            print("AssertionError : Chaque idée doit contenir au maximum 44 caractères\n")
    else: 
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(title)
        for i in range(len(bullet_points)):
            if(i == 0):
                slide.shapes.placeholders[1].text = bullet_points[0]
            else:
                bullet_point = slide.shapes.placeholders[1].text_frame.add_paragraph()
                bullet_point.text = str(bullet_points[i])
                bullet_point.level = levels[i]


def picture_slide(presentation, title = '', image='hello/static/images/python.png', left = 3, top = 2):
    try:
        error = None
        if not type(presentation) is type(Presentation()):
            error = 1
            raise TypeError
        allowed_format = ['.png', '.jpg', '.jpeg']
        allowed = False
        for format in allowed_format:
            if str(image).lower().endswith(format):
                allowed = True
                break
        if not allowed:
            error = 2
            raise TypeError
    except TypeError:
        if error == 1:
            print("TypeError : {} doit être un objet Presentation()\n".format(presentation))
        elif error == 2:
            print("TypeError : L'image est une chaîne de caractères avec le format : {}\n".format(allowed_format))
    else:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = str(title)
        slide.shapes.add_picture(image, Inches(float(left)), Inches(float(top)))


def graph_slide(presentation, title = '', graph_title = '', x_axis = [], y_axis = None, left = 2, top = 2, width = 6, height = 4):
    if(not type(y_axis) is type(array([]))):
        y_axis = ones(len(x_axis))
        print("Attention, vous n'avez pas renseigner les valeurs en ordonnée sous forme d'un un numpy.array. Vérifiez bien\n")
    try:
        error = None
        if not type(presentation) is type(Presentation()):
            error = 1
            raise TypeError
        if not type(x_axis) is list:
            error = 2
            raise TypeError
        if len(x_axis) != len(y_axis):
            error = 1
            raise ValueError
        if len(x_axis) != len(y_axis):
            error = 1
            raise ValueError
    except TypeError:
        if error == 1:
            print("TypeError : {} doit être un objet Presentation()\n".format(presentation))
        elif error == 2:
            print("TypeError : Veuillez renseigner une liste contenant les valeurs en abscisses")
            print("[valeur1, valeur2]\n")
    except ValueError:
        print("ValueError : Les listes des idées et des niveaux doivent être de même taille\n")
    else:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = str(title)
        graph_info = CategoryChartData()
        for i in range(len(x_axis)):
            x_axis[i]= str(x_axis[i])
        graph_info.categories = x_axis
        graph_info.add_series(graph_title, y_axis)
        graph_frame = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(float(left)),
                                            Inches(float(top)), Inches(float(width)), Inches(float(height)), graph_info)
        return graph_frame


def table_slide(presentation, title = '', data = [[]], rows = 3, columns = 4, left = 1, top = 2, width = 8, height = 3):
    try:
        error = None
        if not type(presentation) is type(Presentation()):
            error = 1
            raise TypeError
        if not type(data) is list:
            error = 2
            raise TypeError
        if not type(rows) is int or rows <= 0 or not type(columns) is int or columns <= 0:
            error = 2
            raise AssertionError
        for d in data:
            if not type(d) is list:
                error = 2
                raise TypeError
            if len(d) != columns :
                error = 1
                raise AssertionError
        if len(data) != rows :
            error = 1
            raise AssertionError  
    except TypeError:
        if(error == 1):
            print("TypeError : {} doit être un objet Presentation()\n".format(presentation))
        elif(error == 2):
            print("TypeError : Veuillez renseigner les données sous forme de matrice :")
            print("[['première_donnée', 'autre_donnée'], ['donnée_autre_ligne', ...']]\n")
    except AssertionError:
        if(error == 1):
            print("AssertionError : Veuillez renseigner les données au complet")
            print("sans renseigner plus d'information qu'il y a de colonnes et de lignes\n")
        elif(error == 2):
            print("AssertionError : Le nombre de lignes et de colonnes doivent êtres des entiers positifs non nuls\n")
    else:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = str(title)
        table_frame = slide.shapes.add_table(rows, columns, Inches(float(left)), Inches(float(top)),
                                            Inches(float(width)), Inches(float(height)))
        table = table_frame.table
        for i in range (columns):
            for j in range (rows):
                cell = table.cell(j, i)
                cell.text = str(data[j][i])

def link_slide(presentation, title = '', link_text = 'lien Google', url_link = 'https://www.google.fr/'):
    try:
        error = None
        if not type(presentation) is type(Presentation()):
            error = 1
            raise TypeError
        allowed_protocol = ['https://', 'http://']
        allowed = False
        for protocol in allowed_protocol:
            if str(url_link).lower().startswith(protocol):
                allowed = True
                break
        if not allowed:
            error = 2
            raise TypeError
    except TypeError:
        if(error == 1):
            print("TypeError: {} doit être un objet Presentation()\n".format(presentation))
        elif(error == 2):
            print("TypeError: Le lien est une chaîne de caractères utilisant les protocoles {}\n".format(allowed_protocol))
    else:
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        slide.shapes.title.text = str(title)
        link = slide.shapes.placeholders[1].text_frame.paragraphs[0].add_run()
        link.text = str(link_text)
        link.hyperlink.address = url_link